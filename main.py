import os
from pptx import Presentation

PPT_FILE_PATH = "data/서울관광명소.pptx"
OUTPUT_FOLDER = "output"

def extract_text_from_ppt(ppt_path):
    """PPT 파일에서 텍스트를 추출하는 함수"""
    
    # 1. 파일 존재 여부 확인
    if not os.path.exists(ppt_path):
        return f"오류: '{ppt_path}' 파일을 찾을 수 없습니다. 경로를 확인해주세요."

    try:
        prs = Presentation(ppt_path)
        full_text = [] # 추출한 텍스트를 저장할 리스트

        # 2. 모든 슬라이드를 순회하며 텍스트 추출
        for i, slide in enumerate(prs.slides, start=1):
            full_text.append(f"========== 슬라이드 {i} ==========\n")

            # 도형(shape) 안의 텍스트 추출
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    full_text.append(shape.text)

                # 표(table) 안의 텍스트 추출
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame.text.strip():
                                full_text.append(cell.text_frame.text)
            
            # 슬라이드 노트 텍스트 추출
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                full_text.append("\n--- 노트 ---")
                full_text.append(slide.notes_slide.notes_text_frame.text)
            
            full_text.append("\n") # 슬라이드 간 구분을 위한 공백

        return "\n".join(full_text)

    except Exception as e:
        return f"오류: 파일을 처리하는 중 문제가 발생했습니다.\n{e}"


def save_text_to_file(text_content, original_file_path, output_dir):
    """추출된 텍스트를 파일로 저장하는 함수"""
    
    # 3. 결과 저장 폴더 생성
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        print(f"'{output_dir}' 폴더를 생성했습니다.")

    # 4. 저장할 파일명 생성 및 저장
    base_name = os.path.basename(original_file_path)
    file_name_without_ext = os.path.splitext(base_name)[0]
    output_path = os.path.join(output_dir, f"{file_name_without_ext}.txt")
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(text_content)
        
    return output_path

# --- 메인 코드 실행 --- #
if __name__ == "__main__":
    extracted_content = extract_text_from_ppt(PPT_FILE_PATH)
    
    # 오류 없이 텍스트가 추출되었는지 확인
    if not extracted_content.startswith("오류:"):
        saved_file = save_text_to_file(extracted_content, PPT_FILE_PATH, OUTPUT_FOLDER)
        print(f"✅ 텍스트 추출 완료! >> {saved_file}")
    else:
        # 오류 메시지 출력
        print(f"❌ {extracted_content}")